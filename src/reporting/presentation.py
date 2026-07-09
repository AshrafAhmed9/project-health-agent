from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from typing import Dict, Any

# Brand color definitions
PRIMARY_DARK = RGBColor(0x1B, 0x2A, 0x4A)  # Dark Navy (#1B2A4A)
PRIMARY_MEDIUM = RGBColor(0x2D, 0x4A, 0x7A)  # Medium Blue (#2D4A7A)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)  # Bright Blue (#0096D6)
RED_COLOR = RGBColor(0xE8, 0x3E, 0x3E)  # Red
AMBER_COLOR = RGBColor(0xF5, 0xA6, 0x23)  # Amber
GREEN_COLOR = RGBColor(0x2E, 0xCC, 0x71)  # Green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)


class PPTXGenerator:
    """
    Generates a professional 16:9 executive presentation deck using python-pptx.
    """

    def __init__(self):
        self.prs = Presentation()
        # Set 16:9 Widescreen dimensions
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _add_slide_header(self, slide, title: str, subtitle: str):
        """
        Helper to add a standard structured header bar to each slide.
        """
        # Header shape/bar
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY_DARK
        shape.line.color.rgb = PRIMARY_DARK

        # Title text box
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.9)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        # Slide Title
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(24)
        p1.font.color.rgb = WHITE
        p1.font.name = "Calibri"

        # Slide Subtitle
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.italic = True
        p2.font.size = Pt(12)
        p2.font.color.rgb = ACCENT_BLUE
        p2.font.name = "Calibri"

    def _add_bullet_list(
        self, slide, bullets: list, left: float, top: float, width: float, height: float
    ):
        """
        Helper to add a text box with custom formatted bullet points.
        """
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for idx, bullet_text in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = "• " + bullet_text
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "Calibri"
            p.space_after = Pt(10)

    def _add_speaker_notes(self, slide, notes: str):
        """
        Helper to write text to the slide notes section.
        """
        slide.notes_slide.notes_text_frame.text = notes

    def build_deck(
        self,
        slide_data: Dict[str, Any],
        output_path: Path,
        trend_chart_path: Path,
        resource_chart_path: Path,
    ):
        """
        Generates all 7 slides with data, text layouts, charts, and formats notes.
        """
        blank_layout = self.prs.slide_layouts[6]  # Completely blank slide layout

        # ----------------------------------------------------
        # SLIDE 1: Portfolio Health Overview (Dashboard layout)
        # ----------------------------------------------------
        slide1 = self.prs.slides.add_slide(blank_layout)
        data1 = slide_data.get("slide1", {})
        self._add_slide_header(
            slide1,
            data1.get("title", "Portfolio Health Overview"),
            data1.get("subtitle", ""),
        )
        self._add_bullet_list(slide1, data1.get("bullets", []), 0.5, 1.4, 7.0, 5.0)

        # Add Card for Project S2P (Titan)
        card1 = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.0),
            Inches(1.5),
            Inches(4.8),
            Inches(2.3),
        )
        card1.fill.solid()
        card1.fill.fore_color.rgb = LIGHT_BG
        card1.line.color.rgb = PRIMARY_MEDIUM
        tf1 = card1.text_frame
        tf1.word_wrap = True
        p_c1 = tf1.paragraphs[0]
        p_c1.text = "Project A: S2P Titan Implementation"
        p_c1.font.bold = True
        p_c1.font.size = Pt(16)
        p_c1.font.color.rgb = PRIMARY_DARK
        p_c1_2 = tf1.add_paragraph()
        p_c1_2.text = "Overall Status: 🔴 RED\nCompletion: 71%\nSchedule Index (SPI): 0.77\nUnassigned tasks: 67.5%"
        p_c1_2.font.size = Pt(12)
        p_c1_2.font.color.rgb = DARK_TEXT

        # Add Card for Project Plan B (UniSan)
        card2 = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.0),
            Inches(4.2),
            Inches(4.8),
            Inches(2.3),
        )
        card2.fill.solid()
        card2.fill.fore_color.rgb = LIGHT_BG
        card2.line.color.rgb = PRIMARY_MEDIUM
        tf2 = card2.text_frame
        tf2.word_wrap = True
        p_c2 = tf2.paragraphs[0]
        p_c2.text = "Project B: S2P UniSan Implementation"
        p_c2.font.bold = True
        p_c2.font.size = Pt(16)
        p_c2.font.color.rgb = PRIMARY_DARK
        p_c2_2 = tf2.add_paragraph()
        p_c2_2.text = "Overall Status: 🔴 RED\nCompletion: 44%\nSchedule Index (SPI): 0.79\nUnassigned tasks: 35.4%"
        p_c2_2.font.size = Pt(12)
        p_c2_2.font.color.rgb = DARK_TEXT

        self._add_speaker_notes(slide1, data1.get("speaker_notes", ""))

        # ----------------------------------------------------
        # SLIDE 2: Portfolio Trajectory Trends (Trend line chart embedded)
        # ----------------------------------------------------
        slide2 = self.prs.slides.add_slide(blank_layout)
        data2 = slide_data.get("slide2", {})
        self._add_slide_header(
            slide2,
            data2.get("title", "Portfolio Trajectory Trends"),
            data2.get("subtitle", ""),
        )
        self._add_bullet_list(slide2, data2.get("bullets", []), 0.5, 1.4, 6.0, 5.0)

        # Insert trend chart image
        if trend_chart_path.exists():
            slide2.shapes.add_picture(
                str(trend_chart_path),
                Inches(7.0),
                Inches(1.5),
                Inches(5.8),
                Inches(4.8),
            )

        self._add_speaker_notes(slide2, data2.get("speaker_notes", ""))

        # ----------------------------------------------------
        # SLIDE 3: Critical Risk Matrix (Risk visual grid + text description)
        # ----------------------------------------------------
        slide3 = self.prs.slides.add_slide(blank_layout)
        data3 = slide_data.get("slide3", {})
        self._add_slide_header(
            slide3,
            data3.get("title", "Critical Risk Matrix"),
            data3.get("subtitle", ""),
        )
        self._add_bullet_list(slide3, data3.get("bullets", []), 0.5, 1.4, 7.5, 5.0)

        # Generate simple visual Risk Matrix card representation on slide
        matrix_bg = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8.5), Inches(1.5), Inches(4.3), Inches(4.5)
        )
        matrix_bg.fill.solid()
        matrix_bg.fill.fore_color.rgb = LIGHT_BG
        matrix_bg.line.color.rgb = PRIMARY_DARK
        tf3_m = matrix_bg.text_frame
        tf3_m.word_wrap = True
        p_m = tf3_m.paragraphs[0]
        p_m.text = "CRITICAL RISK HEATMAP"
        p_m.font.bold = True
        p_m.font.size = Pt(14)
        p_m.font.color.rgb = PRIMARY_DARK
        p_m_2 = tf3_m.add_paragraph()
        p_m_2.text = "\n🔥 [HIGH IMPACT / HIGH LIKELIHOOD]\n- Phase 2 P2P Schedule Slip (-81d)\n- S2P Hypercare Stall (6% complete)\n- Titan Resource Deficit (67.5% missing)\n\n⚡ [HIGH IMPACT / MED LIKELIHOOD]\n- Plan B Training Phase I Slip (+17d)\n\n⚠️ [MED IMPACT / MED LIKELIHOOD]\n- Plan B Config Documentation delay"
        p_m_2.font.size = Pt(11)
        p_m_2.font.color.rgb = DARK_TEXT

        self._add_speaker_notes(slide3, data3.get("speaker_notes", ""))

        # ----------------------------------------------------
        # SLIDE 4: Schedule Performance Deep-Dive
        # ----------------------------------------------------
        slide4 = self.prs.slides.add_slide(blank_layout)
        data4 = slide_data.get("slide4", {})
        self._add_slide_header(
            slide4,
            data4.get("title", "Schedule Performance Deep-Dive"),
            data4.get("subtitle", ""),
        )
        self._add_bullet_list(slide4, data4.get("bullets", []), 0.5, 1.4, 12.0, 5.0)
        self._add_speaker_notes(slide4, data4.get("speaker_notes", ""))

        # ----------------------------------------------------
        # SLIDE 5: Resource & Dependency Analysis (Resource gap chart embedded)
        # ----------------------------------------------------
        slide5 = self.prs.slides.add_slide(blank_layout)
        data5 = slide_data.get("slide5", {})
        self._add_slide_header(
            slide5,
            data5.get("title", "Resource & Dependency Analysis"),
            data5.get("subtitle", ""),
        )
        self._add_bullet_list(slide5, data5.get("bullets", []), 0.5, 1.4, 6.0, 5.0)

        # Insert resource chart image
        if resource_chart_path.exists():
            slide5.shapes.add_picture(
                str(resource_chart_path),
                Inches(6.8),
                Inches(1.6),
                Inches(6.0),
                Inches(3.8),
            )

        self._add_speaker_notes(slide5, data5.get("speaker_notes", ""))

        # ----------------------------------------------------
        # SLIDE 6: Strategic Recommendations (Action plan grid)
        # ----------------------------------------------------
        slide6 = self.prs.slides.add_slide(blank_layout)
        data6 = slide_data.get("slide6", {})
        self._add_slide_header(
            slide6,
            data6.get("title", "Strategic Recommendations"),
            data6.get("subtitle", ""),
        )
        self._add_bullet_list(slide6, data6.get("bullets", []), 0.5, 1.4, 12.333, 5.0)
        self._add_speaker_notes(slide6, data6.get("speaker_notes", ""))

        # ----------------------------------------------------
        # SLIDE 7: 30-Day Outlook & Forecast
        # ----------------------------------------------------
        slide7 = self.prs.slides.add_slide(blank_layout)
        data7 = slide_data.get("slide7", {})
        self._add_slide_header(
            slide7,
            data7.get("title", "30-Day Outlook & Forecast"),
            data7.get("subtitle", ""),
        )
        self._add_bullet_list(slide7, data7.get("bullets", []), 0.5, 1.4, 12.333, 5.0)
        self._add_speaker_notes(slide7, data7.get("speaker_notes", ""))

        # Save presentation
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        print(f"🖥️ Executive presentation saved to {output_path}")
        return output_path
